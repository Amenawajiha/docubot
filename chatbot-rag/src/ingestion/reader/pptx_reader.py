"""
PowerPoint Reader - Extract text from presentation files.

Supports:
- Modern PowerPoint (.pptx) - Office Open XML format
- Extracts text from slides, speaker notes, and comments
- Preserves slide order and structure

Output format:
- Slide-by-slide text with clear separation
- Includes titles, body text, tables, and notes
"""

import logging
from io import BytesIO

from .base_reader import BaseReader

logger = logging.getLogger(__name__)


class PPTXReader(BaseReader):
    """
    Extract text from PowerPoint presentation files.
    
    Extracts all readable content:
    - Slide titles
    - Body text and bullet points
    - Tables
    - Speaker notes
    - Comments (optional)
    """
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]
    
    def extract(self, file_bytes: bytes) -> str:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_bytes: Raw PPTX bytes
            
        Returns:
            Slide-by-slide text content
            
        Raises:
            ValueError: If file is empty or contains no text
            RuntimeError: If parsing fails
        """
        if not file_bytes:
            raise ValueError("Empty PowerPoint file provided")
        
        # Lazy import
        try:
            from pptx import Presentation
        except ImportError as e:
            raise RuntimeError(
                "python-pptx is required for PowerPoint support. "
                "Install with: pip install python-pptx"
            ) from e
        
        try:
            # Load presentation
            prs = Presentation(BytesIO(file_bytes))
            
            all_slides = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = self._extract_slide(slide, slide_num)
                
                if slide_text.strip():
                    all_slides.append(slide_text)
            
            if not all_slides:
                raise ValueError("No text content found in PowerPoint")
            
            # Combine slides with clear separation
            content = "\n\n---\n\n".join(all_slides)

            try: 
                from src.ingestion.utils.ocr import run_ocr
            except ImportError:
                run_ocr = None

            ocr_texts = []
            for image_bytes in self._extract_images(file_bytes, Presentation):
                if run_ocr:
                    try:
                        ocr_result = run_ocr(image_bytes)
                        if ocr_result.strip():
                            ocr_texts.append(ocr_result.strip())
                    except Exception as e:
                        logger.warning(f"OCR failed for an image: {e}")

            if ocr_texts:
                content += "\n\n[Image OCR]\n" + "\n\n".join(ocr_texts)
            
            logger.info(
                f"PowerPoint extraction complete: {len(prs.slides)} slides, "
                f"{len(content)} characters"
            )
            return content
        
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise RuntimeError(f"Failed to read PowerPoint: {e}") from e
    
    def _extract_slide(self, slide, slide_num: int) -> str:
        """
        Extract all text from a single slide.
        
        Args:
            slide: python-pptx Slide object
            slide_num: Slide number (1-indexed)
            
        Returns:
            Formatted slide text
        """
        parts = [f"Slide {slide_num}"]
        
        # Extract text from all shapes
        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            
            # Tables
            if hasattr(shape, "table"):
                table_text = self._extract_table(shape.table)
                if table_text:
                    parts.append(table_text)
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes: {notes}]")
        
        return "\n\n".join(parts)
    
    @staticmethod
    def _extract_table(table) -> str:
        """
        Extract text from a table shape.
        
        Args:
            table: python-pptx Table object
            
        Returns:
            Table text in natural language format
        """
        rows_text = []
        
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    cells.append(cell_text)
            
            if cells:
                rows_text.append(" | ".join(cells))
        
        return "\n".join(rows_text) if rows_text else ""
    
    def _extract_images(self, file_bytes: bytes, Presentation):
        prs = Presentation(BytesIO(file_bytes))
        images = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    try:
                        image_bytes = shape.image.blob
                        images.append(image_bytes)
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide: {e}")
        return images


class PPTReader(PPTXReader):
    """
    Legacy PowerPoint reader (.ppt files).
    
    Note: python-pptx only supports .pptx (Office Open XML format).
    For .ppt (binary format), consider using LibreOffice conversion
    or the unoconv library as a preprocessing step.
    """
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".ppt"]
    
    def extract(self, file_bytes: bytes) -> str:
        """
        Extract text from legacy PowerPoint file.
        
        This is a placeholder - .ppt support requires additional libraries.
        """
        raise RuntimeError(
            "Legacy .ppt format is not directly supported. "
            "Please convert to .pptx format first, or use LibreOffice/unoconv "
            "for batch conversion. python-pptx only supports .pptx files."
        )